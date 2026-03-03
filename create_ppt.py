import os
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    import subprocess
    import sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # スライド1: タイトル
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "システムフロー可視化・詳細確認ツール\n（Helios Flow）打合せ"
    subtitle.text = "目的：工事工程の理解・進捗確認・必要ファイル参照を一元化し、\n質問/確認コストを削減する（DX化 / AIX）"

    # スライド2: ゴール
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. 本日の打合せゴール（決めること）"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "システムテスト観点の合意\n（開発部テストで何を見るか）"
    tf.add_paragraph().text = "修正対象の優先順位\n（リリースまでに必ず直す/後回しを切る）"
    tf.add_paragraph().text = "運用イメージの合意\n（誰が、いつ、何を更新するか）"

    # スライド3: 背景
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. 背景：なぜこのツールが必要か（制作意図）"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "工事工程が複雑で、担当者が「次に何をすべきか」を都度確認している状態。"
    tf.add_paragraph().text = "結果として、特定担当に質問が集中し、業務が停滞する。"
    p = tf.add_paragraph()
    p.text = "→ これを「フロー図UIで自己解決できる状態」にする。"
    p.font.bold = True

    # スライド4: 課題定義 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. 課題定義（現状のボトルネック） - 前半"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "3.1 契約書ベースの表管理の限界"
    p = tf.add_paragraph()
    p.text = "契約書に基づく表で進捗管理しているが、詳細の記載が難しく、内容理解が進まない"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "工程の「前後関係」「次にやるべきこと」が表から読み取りづらい"
    p.level = 1
    tf.add_paragraph().text = "3.2 ANDPADの工程管理の限界"
    p = tf.add_paragraph()
    p.text = "基本はバーチャートであり、フロー（分岐/依存）を表現しづらい"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "編集権限が案件監視者に限定され、修正は担当者が都度依頼・反映が遅れる"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "ログイン後にバーチャートを開き、他工程との調整をしながら操作する必要があり、運用負荷が高い"
    p.level = 1

    # スライド5: 課題定義 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. 課題定義（現状のボトルネック） - 後半"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "3.3 必要ファイル/確認ファイルの分散"
    p = tf.add_paragraph()
    p.text = "必要なファイルが一元管理できておらず、参照経路が属人化している"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "工程ごとの「必要資料」「確認すべき資料」を探す時間が発生している"
    p.level = 1
    tf.add_paragraph().text = "3.4 工程内容が複雑で、都度確認が必要"
    p = tf.add_paragraph()
    p.text = "対応内容が複雑で「次の一手」が不明確になり、質問が増える"
    p.level = 1

    # スライド6: ゴール
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4. アプリケーションのゴール（到達状態）"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Helios工程をフロー図化し、工程の前後関係・分岐を視覚的に理解できる"
    tf.add_paragraph().text = "フロー図上で進捗（未着手/進行中/完了）を入力し、状態を可視化できる"
    tf.add_paragraph().text = "遅延フローを可視化し、遅延の早期検知・解消につなげる"
    tf.add_paragraph().text = "開発担当者以外でも、フローを見て進捗確認できる（質問を減らす）"
    tf.add_paragraph().text = "工程（ノード）から「必要ファイル」を抽出し、迷わず参照できる"
    tf.add_paragraph().text = "工程（ノード）から「必要情報（手順/注意点/提出先）」を抽出できる"

    # スライド7: 現在の状況
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "5. 現在の状況（進捗）"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "システムフロー図の作成：完了"
    tf.add_paragraph().text = "フロー可視化（アニメーション等）：実装済み"
    tf.add_paragraph().text = "本日：業務展開（開発部によるシステムテスト開始）"
    tf.add_paragraph().text = "【現在地】"
    p = tf.add_paragraph()
    p.text = "課題抽出 → 要件定義 → システム開発 → 【システムテスト(いまここ)】 → 修正 → リリース"
    p.level = 1

    # スライド8: 今後の進め方・タスク
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "6. 今後の進め方 / 7. タスク"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "【今後のスケジュール】"
    p = tf.add_paragraph()
    p.text = "本日〜1週間程度：開発部へシステムテスト実施"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "テスト結果を踏まえ、修正要件の打合せ（優先順位付け含む）"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "修正対応 → リリース"
    p.level = 1
    tf.add_paragraph().text = "【今後のタスク担当】"
    p = tf.add_paragraph()
    p.text = "若林担当：ノードテキスト検討、関連資料の作成/整理"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "山崎担当：ノード内容修正（文言/説明/リンク整理）"
    p.level = 1

    # スライド9: 本日の確認事項
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "8. 今日の確認事項（会議の論点）"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "システムテストの観点"
    p = tf.add_paragraph()
    p.text = "（UI操作、フロー理解、検索/抽出、進捗入力、遅延表示）"
    p.level = 1
    tf.add_paragraph().text = "修正優先度"
    p = tf.add_paragraph()
    p.text = "（必須 / できれば / 後回し）"
    p.level = 1
    tf.add_paragraph().text = "運用について"
    p = tf.add_paragraph()
    p.text = "（誰がノード情報・リンク・ファイルを更新するか）"
    p.level = 1

    # 保存
    prs.save('HeliosFlow_Meeting.pptx')
    print("Presentation saved as HeliosFlow_Meeting.pptx")

if __name__ == '__main__':
    create_presentation()
